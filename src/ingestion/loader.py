import os
import io
import zipfile
import xml.etree.ElementTree as ET
from typing import List, Set, Optional
import pymupdf as fitz
import numpy as np
from langchain_core.documents import Document
from src.chunking.chunker import split_documents
from src.embeddings.embedder import get_embedding_model
from src.vectordb.vector_store import get_vector_store, VectorStoreManager
from src.utils.helpers import load_app_config


def is_image_pdf(pdf_doc: fitz.Document) -> bool:
    """Check if PDF is primarily scanned images without selectable text."""
    total_text_len = 0
    total_images = 0
    short_text_pages = 0
    for page in pdf_doc:
        page_text_len = len(page.get_text().strip())
        total_text_len += page_text_len
        total_images += len(page.get_images())
        if page_text_len < 30:
            short_text_pages += 1

    page_count = max(len(pdf_doc), 1)
    mostly_short_pages = short_text_pages / page_count >= 0.5
    low_text_density = total_text_len / page_count < 30
    return total_images > 0 and (low_text_density or mostly_short_pages)


def _ocr_page(page: fitz.Page) -> Optional[str]:
    """Extract readable text from a scanned page using the bundled OCR backend."""
    pix = page.get_pixmap(dpi=180, alpha=False)
    image = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)

    try:
        from rapidocr_onnxruntime import RapidOCR

        result, _ = RapidOCR()(image)
        if result:
            lines = [
                item[1].strip()
                for item in result
                if len(item) > 2 and float(item[2]) >= 0.35 and item[1].strip()
            ]
            text = "\n".join(lines)
            if text.strip():
                return text
    except Exception:
        pass

    try:
        import pytesseract
        from PIL import Image

        tesseract_default = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
        if os.path.exists(tesseract_default):
            pytesseract.pytesseract.tesseract_cmd = tesseract_default
        return pytesseract.image_to_string(Image.fromarray(image)).strip()
    except Exception:
        return None


def load_single_pdf(file_path: str) -> List[Document]:
    """Extract text from PDF file with automatic OCR fallback for scanned pages."""
    docs = []
    pdf = fitz.open(file_path)

    if is_image_pdf(pdf):
        print(f"  [OCR SCANNING] {os.path.basename(file_path)} ({len(pdf)} pages)")
        for page_num, page in enumerate(pdf):
            print(f"    -> Page {page_num + 1}/{len(pdf)}", end="\r")
            text = _ocr_page(page)
            if text and len(text.strip()) >= 30:
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "page": page_num + 1, "type": "ocr_scanned"}
                ))
        print()
    else:
        for page_num, page in enumerate(pdf):
            text = page.get_text()
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "page": page_num + 1, "type": "digital_pdf"}
                ))
    pdf.close()
    return docs


def load_single_docx(file_path: str) -> List[Document]:
    """Extract DOCX text, tolerating corrupt embedded images or other media."""
    def load_xml_text() -> str:
        namespace = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        qualified_text = f"{{{namespace}}}t"
        with zipfile.ZipFile(file_path) as archive:
            document_xml = archive.read("word/document.xml")
        root = ET.fromstring(document_xml)
        blocks = []
        for element in root.iter():
            if element.tag == f"{{{namespace}}}p":
                text = "".join(node.text or "" for node in element.iter(qualified_text)).strip()
                if text:
                    blocks.append(text)
        return "\n".join(blocks)

    try:
        import docx
        doc = docx.Document(file_path)
        full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
        if full_text.strip():
            return [Document(page_content=full_text, metadata={"source": file_path, "type": "docx"})]
    except Exception as e:
        try:
            print(f"  [DOCX media warning] {os.path.basename(file_path)}: {e}")
        except Exception:
            pass

        try:
            recovered_text = load_xml_text()
            if recovered_text.strip():
                try:
                    print(f"  [DOCX XML recovered] {os.path.basename(file_path)}")
                except Exception:
                    pass
                return [Document(
                    page_content=recovered_text,
                    metadata={"source": file_path, "type": "docx_xml_recovered"},
                )]
        except Exception as recovery_error:
            try:
                print(f"  [DOCX XML recovery failed] {os.path.basename(file_path)}: {recovery_error}")
            except Exception:
                pass

    return []


def load_single_pptx(file_path: str) -> List[Document]:
    """Extract text from PPTX slides."""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        docs = []
        for idx, slide in enumerate(prs.slides):
            text = "\n".join([
                shape.text for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ])
            if text.strip():
                docs.append(Document(
                    page_content=text,
                    metadata={"source": file_path, "slide": idx + 1, "type": "pptx"}
                ))
        return docs
    except Exception as e:
        print(f"  ⚠️ Skipping unreadable PPTX {os.path.basename(file_path)}: {e}")
    return []


def load_document(file_path: str) -> List[Document]:
    """Universal document loader dispatching by extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return load_single_pdf(file_path)
    elif ext == ".docx":
        return load_single_docx(file_path)
    elif ext == ".pptx":
        return load_single_pptx(file_path)
    return []


def process_directory_incrementally(
    root_path: str = "./PDF_Data",
    subject_filter: Optional[str] = None,
    rebuild: bool = False,
):
    """
    Incrementally ingest documents from folder into the ChromaDB vector database.
    Skips files that are already indexed in ChromaDB.
    """
    cfg = load_app_config()
    db_manager = get_vector_store()
    embedder = get_embedding_model()

    if rebuild:
        deleted = db_manager.delete_records(subject=subject_filter)
        scope = subject_filter or "all subjects"
        print(f"Rebuild requested for {scope}: removed {deleted} vector records.")

    indexed_sources = set() if rebuild else db_manager.get_indexed_sources()
    print(f"Already indexed files in database: {len(indexed_sources)}")

    supported_extensions = {".pdf", ".docx", ".pptx"}
    all_files = []
    for root, dirs, files in os.walk(root_path):
        for file in files:
            ext = os.path.splitext(file)[1].lower()
            if ext in supported_extensions:
                file_path = os.path.join(root, file)
                relative_parts = os.path.relpath(file_path, root_path).replace("\\", "/").split("/")
                if subject_filter and (not relative_parts or relative_parts[0] != subject_filter):
                    continue
                all_files.append(file_path)

    total_files = len(all_files)
    print(f"Found {total_files} total supported documents in '{root_path}'.\n")

    for idx, file_path in enumerate(all_files, start=1):
        normalized_path = file_path.replace("\\", "/")
        if file_path in indexed_sources or normalized_path in indexed_sources:
            print(f"[{idx}/{total_files}] ⏩ SKIP (Already Indexed): {os.path.basename(file_path)}")
            continue

        print(f"[{idx}/{total_files}] 🔄 Processing: {os.path.basename(file_path)}")
        try:
            docs = load_document(file_path)
            if docs:
                chunks = split_documents(docs, root_data_dir=root_path)
                if chunks:
                    texts = [c.page_content for c in chunks]
                    embeddings = embedder.encode(texts)
                    db_manager.add_documents(chunks, embeddings)
                    print(f"  ✅ Saved {len(chunks)} chunks to vector database.")
            else:
                print(f"  ⚠️ No text extracted from {os.path.basename(file_path)}.")
        except Exception as e:
            print(f"  ❌ Error processing file {file_path}: {e}")

    print("\n🎉 Pipeline Execution Completed! Total chunks in DB:", db_manager.collection.count())
